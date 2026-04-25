"""Solar Sizing PowerPoint proposal generator — Rincol Tech Solutions.

Generates a 9-slide widescreen (16:9) proposal from sizing data.
Slides: Cover | Why Solar | Energy Profile | System Sizing |
        Equipment/BoM | Financial Appraisal | Payback Timeline |
        Why Rincol | Call to Action
"""
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ─────────────────────────────────────────────────────────────
DARK_NAVY  = RGBColor(0x1a, 0x1a, 0x2e)
MID_NAVY   = RGBColor(0x16, 0x21, 0x3e)
BLUE       = RGBColor(0x23, 0x6F, 0xC4)
BLUE_LIGHT = RGBColor(0x89, 0xb4, 0xfa)
GREEN      = RGBColor(0x22, 0xc5, 0x5e)
AMBER      = RGBColor(0xf5, 0x9e, 0x0b)
WHITE      = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE  = RGBColor(0xe2, 0xe8, 0xf0)
GREY_TEXT  = RGBColor(0x94, 0xa3, 0xb8)
DIVIDER    = RGBColor(0x2a, 0x2d, 0x45)

W = Inches(13.33)   # 16:9 widescreen width
H = Inches(7.5)     # 16:9 widescreen height


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def _txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _bg(slide):
    """Dark navy full-slide background."""
    _rect(slide, 0, 0, W, H, DARK_NAVY)


def _header_bar(slide, title_text, subtitle_text=''):
    """Blue accent bar at top with title."""
    _rect(slide, 0, 0, W, Inches(1.2), BLUE)
    _txt(slide, title_text, Inches(0.4), Inches(0.15), Inches(12), Inches(0.6),
         size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        _txt(slide, subtitle_text, Inches(0.4), Inches(0.75), Inches(12), Inches(0.4),
             size=14, color=OFF_WHITE)


def _kpi_card(slide, x, y, w, h, label, value, value_color=WHITE):
    """A metric card with a mid-navy background."""
    _rect(slide, x, y, w, h, MID_NAVY)
    _txt(slide, label, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.35),
         size=10, color=GREY_TEXT, align=PP_ALIGN.CENTER)
    _txt(slide, value, x + Inches(0.1), y + Inches(0.4), w - Inches(0.2), Inches(0.55),
         size=18, bold=True, color=value_color, align=PP_ALIGN.CENTER)


def _fmt_ugx(v):
    try:
        return f"UGX {int(v):,}"
    except Exception:
        return str(v)


def _fmt_n(v, dp=1):
    try:
        return f"{float(v):,.{dp}f}"
    except Exception:
        return str(v)


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_cover(prs, s, r):
    slide = _blank(prs)
    _bg(slide)
    # Diagonal blue accent block
    _rect(slide, 0, 0, Inches(6), H, MID_NAVY)
    _rect(slide, 0, 0, Inches(0.25), H, BLUE)
    # Company name
    _txt(slide, "RINCOL TECH SOLUTIONS", Inches(0.5), Inches(1.2), Inches(5), Inches(0.6),
         size=14, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)
    # Main heading
    _txt(slide, "Solar Investment\nProposal", Inches(0.5), Inches(2.0), Inches(5), Inches(1.8),
         size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Client name
    _txt(slide, f"Prepared for: {s['client_name']}", Inches(0.5), Inches(4.0),
         Inches(5.5), Inches(0.5), size=16, color=OFF_WHITE)
    if s.get('client_site'):
        _txt(slide, s['client_site'], Inches(0.5), Inches(4.5), Inches(5.5), Inches(0.4),
             size=12, color=GREY_TEXT, italic=True)
    # Right side stats
    _txt(slide, "System Summary", Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
         size=14, bold=True, color=BLUE_LIGHT)
    stats = [
        ("Daily Load",    f"{_fmt_n(r['total_daily_wh'], 0)} Wh/day"),
        ("Solar Array",   f"{r['panels_recommended']} × {_fmt_n(s['panel_wp'], 0)}Wp panels"),
        ("Battery Bank",  f"{r['total_batteries']} × {_fmt_n(s['battery_ah'], 0)}Ah / {s['battery_type']}"),
        ("Inverter",      f"{s['inverter_kw']} kW / {s['system_voltage']}V"),
        ("Payback",       f"{_fmt_n(r['payback_years'])} years"),
    ]
    for i, (lbl, val) in enumerate(stats):
        y = Inches(2.1) + i * Inches(0.9)
        _rect(slide, Inches(7.0), y, Inches(5.7), Inches(0.75), MID_NAVY)
        _txt(slide, lbl, Inches(7.15), y + Inches(0.08), Inches(2.5), Inches(0.35),
             size=10, color=GREY_TEXT)
        _txt(slide, val, Inches(9.7), y + Inches(0.08), Inches(2.8), Inches(0.35),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def _slide_why_solar(prs, s, r):
    slide = _blank(prs)
    _bg(slide)
    _header_bar(slide, "Why Solar?", "The case for off-grid solar in Uganda")
    points = [
        (BLUE_LIGHT, "Eliminate grid dependency",
         f"{s.get('utility_provider','UEDCL')} tariffs are rising. Solar locks in your energy cost for 25+ years."),
        (GREEN, "Proven payback",
         f"Based on your load, solar pays back in {_fmt_n(r['payback_years'])} years. After that, energy is nearly free."),
        (AMBER, "10-year savings",
         f"Estimated savings of {_fmt_ugx(r['yaka_savings_10yr'])} over 10 years vs grid power."),
        (BLUE_LIGHT, "Reliability",
         "No load shedding. No token shortages. 1 day of battery autonomy even with zero sun."),
        (GREEN, "Low maintenance",
         "Li-ion batteries: zero maintenance for 10 years. Systems designed for Uganda's climate."),
    ]
    for i, (color, title, body) in enumerate(points):
        y = Inches(1.4) + i * Inches(1.1)
        _rect(slide, Inches(0.3), y + Inches(0.05), Inches(0.08), Inches(0.65), color)
        _txt(slide, title, Inches(0.6), y, Inches(12), Inches(0.4),
             size=13, bold=True, color=color)
        _txt(slide, body, Inches(0.6), y + Inches(0.4), Inches(12), Inches(0.55),
             size=11, color=OFF_WHITE)


def _slide_energy_profile(prs, s, r, appliances):
    slide = _blank(prs)
    _bg(slide)
    _header_bar(slide, "Energy Profile", f"Daily load: {_fmt_n(r['total_daily_wh'], 0)} Wh/day  |  Peak: {_fmt_n(r['peak_load_w'], 0)} W")

    # Appliance table
    col_x = [Inches(0.3), Inches(3.5), Inches(5.2), Inches(6.6), Inches(8.0), Inches(9.6), Inches(11.2)]
    headers = ["Appliance", "Power (W)", "Pwr Factor", "Qty", "Hours/day", "Wh/day", "Inc."]
    col_w   = [Inches(3.0), Inches(1.5), Inches(1.3), Inches(1.2), Inches(1.4), Inches(1.5), Inches(1.2)]

    # Header row
    _rect(slide, Inches(0.3), Inches(1.3), Inches(12.8), Inches(0.38), MID_NAVY)
    for j, (hdr, cx) in enumerate(zip(headers, col_x)):
        _txt(slide, hdr, cx, Inches(1.33), col_w[j], Inches(0.35),
             size=9, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)

    # Data rows (max 14 fit on slide)
    incl = [a for a in appliances if a.get('included', True)]
    excl = [a for a in appliances if not a.get('included', True)]
    rows = incl[:14]
    for i, a in enumerate(rows):
        y = Inches(1.7) + i * Inches(0.35)
        if i % 2 == 0:
            _rect(slide, Inches(0.3), y, Inches(12.8), Inches(0.35), RGBColor(0x1e, 0x1e, 0x38))
        vals = [
            a['name'],
            _fmt_n(a['power_w'], 0),
            _fmt_n(a['power_factor'], 2),
            str(a['quantity']),
            _fmt_n(a['hours_per_day'], 1),
            _fmt_n(a.get('daily_wh', 0), 0),
            "Yes",
        ]
        for j, (val, cx) in enumerate(zip(vals, col_x)):
            _txt(slide, val, cx, y + Inches(0.03), col_w[j], Inches(0.32),
                 size=9, color=WHITE, align=PP_ALIGN.LEFT)

    if excl:
        excluded_names = ", ".join(a['name'] for a in excl[:5])
        _txt(slide, f"Excluded from solar: {excluded_names}",
             Inches(0.3), Inches(7.0), Inches(12.8), Inches(0.35),
             size=9, color=GREY_TEXT, italic=True)


def _slide_system_sizing(prs, s, r):
    slide = _blank(prs)
    _bg(slide)
    _header_bar(slide, "System Sizing", "Engineering specification")

    card_w = Inches(3.8)
    card_h = Inches(2.5)

    # Battery card
    bx = Inches(0.4)
    by = Inches(1.5)
    _rect(slide, bx, by, card_w, card_h, MID_NAVY)
    _rect(slide, bx, by, card_w, Inches(0.4), BLUE)
    _txt(slide, "Battery Bank", bx + Inches(0.15), by + Inches(0.05), card_w, Inches(0.35),
         size=13, bold=True, color=WHITE)
    lines = [
        f"Min. Ah required: {_fmt_n(r['battery_ah_min'], 1)} Ah",
        f"Config: {r['batteries_in_series']}S × {r['batteries_in_parallel']}P",
        f"Total: {r['total_batteries']} batteries",
        f"{_fmt_n(s['battery_ah'], 0)}Ah/{_fmt_n(s['battery_voltage'], 0)}V {s['battery_type']}",
    ]
    for i, line in enumerate(lines):
        _txt(slide, line, bx + Inches(0.15), by + Inches(0.5 + i * 0.45),
             card_w - Inches(0.3), Inches(0.42), size=11, color=OFF_WHITE)

    # Solar card
    sx = Inches(4.7)
    sy = Inches(1.5)
    _rect(slide, sx, sy, card_w, card_h, MID_NAVY)
    _rect(slide, sx, sy, card_w, Inches(0.4), BLUE)
    _txt(slide, "Solar Array", sx + Inches(0.15), sy + Inches(0.05), card_w, Inches(0.35),
         size=13, bold=True, color=WHITE)
    override_note = " (voltage override)" if r.get('voltage_override') else ""
    lines = [
        f"Required: {_fmt_n(r['required_wp'], 0)} Wp",
        f"Panels by energy: {r['panels_by_energy']}",
        f"Panels installed: {r['panels_recommended']}{override_note}",
        f"{_fmt_n(s['panel_wp'], 0)}Wp monocrystalline",
    ]
    for i, line in enumerate(lines):
        _txt(slide, line, sx + Inches(0.15), sy + Inches(0.5 + i * 0.45),
             card_w - Inches(0.3), Inches(0.42), size=11, color=OFF_WHITE)

    # Inverter card
    ix = Inches(9.0)
    iy = Inches(1.5)
    _rect(slide, ix, iy, card_w, card_h, MID_NAVY)
    _rect(slide, ix, iy, card_w, Inches(0.4), BLUE)
    _txt(slide, "Inverter & System", ix + Inches(0.15), iy + Inches(0.05), card_w, Inches(0.35),
         size=13, bold=True, color=WHITE)
    lines = [
        f"{s['inverter_kw']} kW / {s['system_voltage']}V Hybrid",
        f"Peak load: {_fmt_n(r['peak_load_w']/1000, 2)} kW",
        f"PSH: {s['peak_sun_hours']} hrs/day",
        f"Annual yield: {_fmt_n(r['annual_yield_kwh'], 0)} kWh/yr",
    ]
    flag = r.get('inverter_flag', '')
    flag_color = AMBER if flag else GREY_TEXT
    for i, line in enumerate(lines):
        _txt(slide, line, ix + Inches(0.15), iy + Inches(0.5 + i * 0.45),
             card_w - Inches(0.3), Inches(0.42), size=11, color=OFF_WHITE)

    if flag:
        _rect(slide, Inches(0.4), Inches(4.3), Inches(12.6), Inches(0.5), RGBColor(0x3a, 0x1a, 0x00))
        _txt(slide, f"⚠  {flag}", Inches(0.6), Inches(4.35), Inches(12.4), Inches(0.42),
             size=10, bold=True, color=AMBER)

    # PSH note
    _txt(slide, f"Peak sun hours: {s['peak_sun_hours']} hrs/day  |  Performance ratio: {s['performance_ratio']}  |  DOD: {s['dod']}",
         Inches(0.4), Inches(5.0), Inches(12.6), Inches(0.35), size=9, color=GREY_TEXT)


def _slide_bom(prs, bom_items):
    slide = _blank(prs)
    _bg(slide)
    total = sum(item.get('total', 0) for item in bom_items)
    _header_bar(slide, "Equipment & Bill of Materials",
                f"Total system cost: {_fmt_ugx(total)}")

    col_x = [Inches(0.3), Inches(7.5), Inches(9.0), Inches(10.2), Inches(11.5)]
    headers = ["Item Description", "UoM", "Qty", "Unit Price (UGX)", "Total (UGX)"]
    col_w   = [Inches(7.0), Inches(1.3), Inches(1.0), Inches(1.3), Inches(1.5)]
    aligns  = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]

    _rect(slide, Inches(0.3), Inches(1.3), Inches(12.8), Inches(0.38), MID_NAVY)
    for j, (hdr, cx) in enumerate(zip(headers, col_x)):
        _txt(slide, hdr, cx, Inches(1.33), col_w[j], Inches(0.35),
             size=9, bold=True, color=BLUE_LIGHT, align=aligns[j])

    for i, item in enumerate(bom_items[:15]):
        y = Inches(1.7) + i * Inches(0.32)
        if i % 2 == 0:
            _rect(slide, Inches(0.3), y, Inches(12.8), Inches(0.32), RGBColor(0x1e, 0x1e, 0x38))
        vals = [
            item['description'],
            item.get('uom', 'pc'),
            _fmt_n(item.get('qty', 1), 0),
            f"{int(item.get('unit_price', 0)):,}",
            f"{int(item.get('total', 0)):,}",
        ]
        for j, (val, cx) in enumerate(zip(vals, col_x)):
            _txt(slide, val, cx, y + Inches(0.02), col_w[j], Inches(0.30),
                 size=9, color=WHITE, align=aligns[j])

    # Total row
    ty = Inches(1.7) + len(bom_items[:15]) * Inches(0.32) + Inches(0.05)
    _rect(slide, Inches(0.3), ty, Inches(12.8), Inches(0.38), BLUE)
    _txt(slide, "TOTAL", Inches(0.4), ty + Inches(0.05), Inches(10.8), Inches(0.32),
         size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    _txt(slide, f"{int(total):,}", Inches(11.5), ty + Inches(0.05), Inches(1.5), Inches(0.32),
         size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def _slide_financial(prs, s, r):
    slide = _blank(prs)
    _bg(slide)
    _header_bar(slide, "10-Year Financial Appraisal (Section C)",
                f"Utility: {s.get('utility_provider','UEDCL')} @ {_fmt_n(s['utility_tariff'], 0)} UGX/kWh")

    maint = r['maintenance_cost_10yr']
    coo   = r['system_cost'] + maint

    kpis = [
        ("System Cost", _fmt_ugx(r['system_cost']), WHITE),
        ("10-yr Maintenance", _fmt_ugx(maint), AMBER if maint > 0 else GREY_TEXT),
        ("Cost of Ownership (10yr)", _fmt_ugx(coo), WHITE),
        ("Solar Cost/kWh", f"{_fmt_ugx(r['solar_cost_per_kwh'])}", BLUE_LIGHT),
        (f"{s.get('utility_provider','Grid')} Cost/kWh", f"UGX {_fmt_n(s['utility_tariff'], 0)}", AMBER),
        ("10-yr Grid Savings", _fmt_ugx(r['yaka_savings_10yr']),
         GREEN if r['yaka_savings_10yr'] > 0 else AMBER),
    ]

    card_w = Inches(3.9)
    card_h = Inches(1.5)
    for i, (lbl, val, color) in enumerate(kpis):
        col = i % 3
        row = i // 3
        cx = Inches(0.4) + col * (card_w + Inches(0.3))
        cy = Inches(1.5) + row * (card_h + Inches(0.25))
        _kpi_card(slide, cx, cy, card_w, card_h, lbl, val, color)

    # Interpretation text
    _txt(slide, (
        f"Annual energy yield: {_fmt_n(r['annual_yield_kwh'], 0)} kWh  |  "
        f"10-year yield: {_fmt_n(r['annual_yield_kwh']*10, 0)} kWh  |  "
        f"Battery: {s['battery_type']}"
        + (f"  |  Battery replacement included in 10-yr maintenance." if maint > 0 else "  |  No battery replacement needed in 10 years (Li-ion).")
    ), Inches(0.4), Inches(5.0), Inches(12.6), Inches(0.6), size=10, color=GREY_TEXT)


def _slide_payback(prs, s, r):
    slide = _blank(prs)
    _bg(slide)
    _header_bar(slide, "Payback Period (Section D)",
                f"Benchmark tariff: {_fmt_n(s['payback_tariff'], 0)} UGX/kWh (fixed)")

    annual_savings = r['annual_yield_kwh'] * s['payback_tariff']

    _kpi_card(slide, Inches(0.4), Inches(1.5), Inches(3.8), Inches(1.4),
              "Annual Energy Yield", f"{_fmt_n(r['annual_yield_kwh'], 0)} kWh/yr", BLUE_LIGHT)
    _kpi_card(slide, Inches(4.5), Inches(1.5), Inches(3.8), Inches(1.4),
              "Annual Savings (grid avoided)", _fmt_ugx(annual_savings), GREEN)
    _kpi_card(slide, Inches(8.6), Inches(1.5), Inches(4.4), Inches(1.4),
              "Payback Period", f"{_fmt_n(r['payback_years'])} years", GREEN)

    # Timeline bar
    _rect(slide, Inches(0.4), Inches(3.2), Inches(12.6), Inches(0.12), MID_NAVY)
    payback_fraction = min(r['payback_years'] / 10.0, 1.0)
    _rect(slide, Inches(0.4), Inches(3.2), Inches(12.6) * payback_fraction, Inches(0.12), GREEN)

    # Year labels
    for yr in range(1, 11):
        x = Inches(0.4) + Inches(12.6) * yr / 10
        _txt(slide, str(yr), x - Inches(0.2), Inches(3.35), Inches(0.4), Inches(0.3),
             size=9, color=GREY_TEXT, align=PP_ALIGN.CENTER)
        _rect(slide, x, Inches(3.15), Inches(0.01), Inches(0.25), DIVIDER)

    # Payback marker
    px = Inches(0.4) + Inches(12.6) * payback_fraction
    _rect(slide, px - Inches(0.015), Inches(3.0), Inches(0.03), Inches(0.4), AMBER)
    _txt(slide, f"Year {_fmt_n(r['payback_years'])}\nPayback", px - Inches(0.5), Inches(2.55),
         Inches(1.0), Inches(0.45), size=9, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

    # Cumulative savings table
    _txt(slide, "Cumulative savings vs grid (10-year projection)",
         Inches(0.4), Inches(3.8), Inches(12.6), Inches(0.35), size=11, bold=True, color=OFF_WHITE)
    for yr in range(1, 11):
        x = Inches(0.4) + (yr - 1) * Inches(1.26)
        cum = annual_savings * yr - r['system_cost']
        color = GREEN if cum >= 0 else AMBER
        _rect(slide, x, Inches(4.2), Inches(1.1), Inches(1.8) * min(abs(cum) / max(annual_savings * 10, 1), 1.0) + Inches(0.1),
              color if cum >= 0 else RGBColor(0x7f, 0x1d, 0x1d))
        _txt(slide, f"Yr {yr}", x, Inches(6.1), Inches(1.1), Inches(0.3),
             size=8, color=GREY_TEXT, align=PP_ALIGN.CENTER)


def _slide_why_rincol(prs):
    slide = _blank(prs)
    _bg(slide)
    _header_bar(slide, "Why Rincol Tech Solutions?")
    points = [
        ("Experienced", "Solar installations across residential and commercial sites in Uganda since 2022."),
        ("Quality components", "We source quality-certified panels, Li-ion and tubular gel batteries, and hybrid inverters."),
        ("Warranty", "12-month installation warranty on all systems. Ongoing maintenance support available."),
        ("Honest sizing", "We size what your load actually needs — not oversized to inflate the bill."),
        ("Local support", "Kampala-based team. Fast response for service calls and spare parts."),
    ]
    for i, (title, body) in enumerate(points):
        y = Inches(1.5) + i * Inches(1.0)
        _rect(slide, Inches(0.3), y + Inches(0.1), Inches(0.06), Inches(0.55), BLUE)
        _txt(slide, title, Inches(0.55), y + Inches(0.05), Inches(12), Inches(0.38),
             size=13, bold=True, color=BLUE_LIGHT)
        _txt(slide, body, Inches(0.55), y + Inches(0.42), Inches(12), Inches(0.45),
             size=11, color=OFF_WHITE)


def _slide_cta(prs, s, r):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, H, MID_NAVY)
    _rect(slide, 0, 0, Inches(0.3), H, BLUE)
    _txt(slide, "Ready to Go Solar?", Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
         size=36, bold=True, color=WHITE)
    _txt(slide, f"This proposal is valid for 30 days.", Inches(0.6), Inches(2.4),
         Inches(12), Inches(0.5), size=14, color=OFF_WHITE)
    _txt(slide, (
        f"System: {r['panels_recommended']} × {_fmt_n(s['panel_wp'], 0)}Wp panels  |  "
        f"{r['total_batteries']} × {_fmt_n(s['battery_ah'], 0)}Ah {s['battery_type']}  |  "
        f"{s['inverter_kw']} kW Inverter"
    ), Inches(0.6), Inches(3.0), Inches(12.4), Inches(0.5), size=13, color=BLUE_LIGHT)
    _txt(slide, _fmt_ugx(r['system_cost']), Inches(0.6), Inches(3.7), Inches(12), Inches(0.8),
         size=40, bold=True, color=GREEN)
    _txt(slide, "Contact us to proceed", Inches(0.6), Inches(4.6), Inches(12), Inches(0.5),
         size=16, bold=True, color=WHITE)
    _txt(slide, "Rincol Tech Solutions  |  Kampala, Uganda", Inches(0.6), Inches(5.2),
         Inches(12), Inches(0.4), size=12, color=GREY_TEXT)


# ── Public entry point ────────────────────────────────────────────────────────

def build_proposal(sizing_row, results, appliances, bom_items):
    """Generate the full 9-slide proposal.

    Args:
        sizing_row: dict — solar_sizings DB row
        results: dict — output of calc_sizing()
        appliances: list — annotated appliances from calc_sizing()
        bom_items: list — output of build_bom()

    Returns:
        bytes — PPTX file content
    """
    s   = dict(sizing_row)
    r   = dict(results)
    prs = _prs()

    _slide_cover(prs, s, r)
    _slide_why_solar(prs, s, r)
    _slide_energy_profile(prs, s, r, appliances)
    _slide_system_sizing(prs, s, r)
    _slide_bom(prs, bom_items)
    _slide_financial(prs, s, r)
    _slide_payback(prs, s, r)
    _slide_why_rincol(prs)
    _slide_cta(prs, s, r)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
